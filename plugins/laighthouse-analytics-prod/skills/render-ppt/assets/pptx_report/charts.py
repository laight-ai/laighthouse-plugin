"""Native OOXML combo-chart injection for python-pptx slides.

python-pptx's add_chart() cannot build bar+line combo charts on dual axes,
so this module builds the chart part XML by hand (the c:chartSpace schema is
shared between docx and pptx), registers it in the package, and appends a
graphicFrame referencing it to the slide's shape tree.

Series fills, gridlines, axis typography, and legend placement are all set
explicitly from theme.py so the chart matches the HTML report's Chart.js
rendering instead of Word/PowerPoint's default Office theme.
"""
from xml.sax.saxutils import escape as xml_escape

from pptx.opc.package import Part
from pptx.opc.packuri import PackURI
from pptx.oxml import parse_xml

import theme

CHART_CONTENT_TYPE = "application/vnd.openxmlformats-officedocument.drawingml.chart+xml"
CHART_RELATIONSHIP_TYPE = (
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/chart"
)

_CHART_NS = (
    'xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart" '
    'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
    'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"'
)

_CAT_AX = 111111111
_BAR_VAL_AX = 222222222
_LINE_VAL_AX = 333333333


def add_combo_chart(slide, package, categories, bar_series, line_series,
                    x, y, cx, cy, title=None):
    partname = _next_chart_partname(package)
    xml = _build_chart_xml(categories, bar_series, line_series, title)
    chart_part = Part(partname, CHART_CONTENT_TYPE, package, xml.encode("utf-8"))
    rId = slide.part.relate_to(chart_part, CHART_RELATIONSHIP_TYPE)
    _append_graphic_frame(slide, rId, x, y, cx, cy)
    return rId


def _next_chart_partname(package):
    existing = [p for p in package.iter_parts()
                if "/ppt/charts/chart" in str(p.partname)]
    return PackURI(f"/ppt/charts/chart{len(existing) + 1}.xml")


def _append_graphic_frame(slide, rId, x, y, cx, cy):
    shape_id = max((sp.shape_id for sp in slide.shapes), default=0) + 1
    xml = f'''<p:graphicFrame xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
        xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
        xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart"
        xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvGraphicFramePr>
    <p:cNvPr id="{shape_id}" name="Chart {shape_id}"/>
    <p:cNvGraphicFramePr/>
    <p:nvPr/>
  </p:nvGraphicFramePr>
  <p:xfrm><a:off x="{int(x)}" y="{int(y)}"/><a:ext cx="{int(cx)}" cy="{int(cy)}"/></p:xfrm>
  <a:graphic>
    <a:graphicData uri="http://schemas.openxmlformats.org/drawingml/2006/chart">
      <c:chart r:id="{rId}"/>
    </a:graphicData>
  </a:graphic>
</p:graphicFrame>'''.encode("utf-8")
    slide.shapes._spTree.append(parse_xml(xml))


# ── style fragments ────────────────────────────────────────────────────

def _solid_fill(hex_color):
    return f'<a:solidFill><a:srgbClr val="{hex_color.lstrip("#").upper()}"/></a:solidFill>'


def _text_props(size_hundredths, color, bold="0"):
    """c:txPr block: font size (in 1/100 pt), color, and Korean-capable face."""
    return (
        '<c:txPr><a:bodyPr/><a:lstStyle/><a:p><a:pPr>'
        f'<a:defRPr sz="{size_hundredths}" b="{bold}">{_solid_fill(color)}'
        f'<a:latin typeface="{theme.FONT}"/><a:ea typeface="{theme.FONT}"/>'
        '</a:defRPr></a:pPr><a:endParaRPr lang="ko-KR"/></a:p></c:txPr>'
    )


def _str_cache(values):
    pts = "".join(f'<c:pt idx="{i}"><c:v>{xml_escape(v)}</c:v></c:pt>' for i, v in enumerate(values))
    return f'<c:strCache><c:ptCount val="{len(values)}"/>{pts}</c:strCache>'


def _num_cache(values):
    pts = "".join(f'<c:pt idx="{i}"><c:v>{v}</c:v></c:pt>' for i, v in enumerate(values))
    return (
        '<c:numCache><c:formatCode>General</c:formatCode>'
        f'<c:ptCount val="{len(values)}"/>{pts}</c:numCache>'
    )


def _col_letter(idx):
    return chr(66 + idx)  # B, C, D, ... (column A is reserved for categories)


# CT_*Ser is an ordered sequence: idx, order, tx, spPr, (marker), cat, val,
# (smooth). PowerPoint refuses charts whose children are out of order.

def _tx_ref(idx, name):
    col = _col_letter(idx)
    return f'<c:tx><c:strRef><c:f>Sheet1!${col}$1</c:f>{_str_cache([name])}</c:strRef></c:tx>'


def _cat_val_refs(idx, categories, values):
    col = _col_letter(idx)
    return f'''
      <c:cat><c:strRef><c:f>Sheet1!$A$2:$A${len(categories) + 1}</c:f>{_str_cache(categories)}</c:strRef></c:cat>
      <c:val><c:numRef><c:f>Sheet1!${col}$2:${col}${len(values) + 1}</c:f>{_num_cache(values)}</c:numRef></c:val>'''


def _bar_series_xml(idx, series, categories):
    color = series.get("color", theme.CHART_BAR_COLORS[idx % len(theme.CHART_BAR_COLORS)])
    return f'''
    <c:ser>
      <c:idx val="{idx}"/>
      <c:order val="{idx}"/>
      {_tx_ref(idx, series["name"])}
      <c:spPr>{_solid_fill(color)}<a:ln><a:noFill/></a:ln></c:spPr>
      {_cat_val_refs(idx, categories, series["values"])}
    </c:ser>'''


def _line_series_xml(idx, series, categories):
    color = series.get("color", theme.CHART_LINE_COLOR)
    fill = _solid_fill(color)
    return f'''
    <c:ser>
      <c:idx val="{idx}"/>
      <c:order val="{idx}"/>
      {_tx_ref(idx, series["name"])}
      <c:spPr><a:ln w="28575" cap="rnd">{fill}</a:ln></c:spPr>
      <c:marker>
        <c:symbol val="circle"/>
        <c:size val="6"/>
        <c:spPr>{fill}<a:ln w="12700">{_solid_fill("FFFFFF")}</a:ln></c:spPr>
      </c:marker>
      {_cat_val_refs(idx, categories, series["values"])}
      <c:smooth val="0"/>
    </c:ser>'''


def _build_chart_xml(categories, bar_series, line_series, title):
    bar_xml = "".join(
        _bar_series_xml(i, s, categories) for i, s in enumerate(bar_series)
    )
    line_xml = _line_series_xml(len(bar_series), line_series, categories)
    axis_text = _text_props(1000, theme.TEXT_MUTED)
    grid_line = (
        '<c:majorGridlines><c:spPr><a:ln w="9525">'
        f'{_solid_fill(theme.CHART_GRID)}</a:ln></c:spPr></c:majorGridlines>'
    )
    no_axis_line = '<c:spPr><a:ln><a:noFill/></a:ln></c:spPr>'
    title_xml = ""
    if title:
        title_xml = (
            "<c:title><c:tx><c:rich><a:bodyPr/><a:p><a:r><a:t>"
            f"{xml_escape(title)}</a:t></a:r></a:p></c:rich></c:tx><c:overlay val=\"0\"/></c:title>"
        )
    return f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<c:chartSpace {_CHART_NS}>
  <c:chart>
    {title_xml}
    <c:autoTitleDeleted val="1"/>
    <c:plotArea>
      <c:layout/>
      <c:barChart>
        <c:barDir val="col"/>
        <c:grouping val="clustered"/>
        <c:varyColors val="0"/>
        {bar_xml}
        <c:gapWidth val="60"/>
        <c:overlap val="-10"/>
        <c:axId val="{_CAT_AX}"/>
        <c:axId val="{_BAR_VAL_AX}"/>
      </c:barChart>
      <c:lineChart>
        <c:grouping val="standard"/>
        <c:varyColors val="0"/>
        {line_xml}
        <c:axId val="{_CAT_AX}"/>
        <c:axId val="{_LINE_VAL_AX}"/>
      </c:lineChart>
      <c:catAx>
        <c:axId val="{_CAT_AX}"/>
        <c:scaling><c:orientation val="minMax"/></c:scaling>
        <c:delete val="0"/>
        <c:axPos val="b"/>
        <c:majorTickMark val="none"/>
        <c:spPr><a:ln w="9525">{_solid_fill(theme.CHART_GRID)}</a:ln></c:spPr>
        {axis_text}
        <c:crossAx val="{_BAR_VAL_AX}"/>
      </c:catAx>
      <c:valAx>
        <c:axId val="{_BAR_VAL_AX}"/>
        <c:scaling><c:orientation val="minMax"/></c:scaling>
        <c:delete val="0"/>
        <c:axPos val="l"/>
        {grid_line}
        <c:numFmt formatCode="#,##0" sourceLinked="0"/>
        <c:majorTickMark val="none"/>
        {no_axis_line}
        {axis_text}
        <c:crossAx val="{_CAT_AX}"/>
      </c:valAx>
      <c:valAx>
        <c:axId val="{_LINE_VAL_AX}"/>
        <c:scaling><c:orientation val="minMax"/></c:scaling>
        <c:delete val="0"/>
        <c:axPos val="r"/>
        <c:numFmt formatCode="#,##0" sourceLinked="0"/>
        <c:majorTickMark val="none"/>
        {no_axis_line}
        {axis_text}
        <c:crossAx val="{_CAT_AX}"/>
        <c:crosses val="max"/>
      </c:valAx>
      <c:spPr><a:noFill/><a:ln><a:noFill/></a:ln></c:spPr>
    </c:plotArea>
    <c:legend>
      <c:legendPos val="b"/>
      <c:overlay val="0"/>
      {_text_props(1100, theme.TEXT_TH)}
    </c:legend>
    <c:plotVisOnly val="1"/>
  </c:chart>
  <c:spPr><a:noFill/><a:ln><a:noFill/></a:ln></c:spPr>
</c:chartSpace>'''
