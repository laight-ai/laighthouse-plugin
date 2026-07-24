from lxml import etree
from pptx import Presentation

import charts
import theme

_C = "{http://schemas.openxmlformats.org/drawingml/2006/chart}"
_A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"

_BARS = [
    {"name": "광고비", "values": [1, 2, 3], "color": "94A3B8"},
    {"name": "전환매출", "values": [4, 5, 6], "color": "93C5FD"},
]
_LINE = {"name": "ROAS(%)", "values": [7, 8, 9], "color": "EF4444"}


def _slide_and_chart():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    charts.add_combo_chart(slide, prs.part.package, ["1월", "2월", "3월"],
                           _BARS, _LINE, x=0, y=0, cx=100000, cy=100000)
    part = next(p for p in prs.part.package.iter_parts()
                if "charts/chart" in str(p.partname))
    return slide, etree.fromstring(part.blob)


def test_chart_part_created_and_frame_on_slide():
    slide, root = _slide_and_chart()
    assert root.tag == f"{_C}chartSpace"
    frames = slide.shapes._spTree.findall(
        "{http://schemas.openxmlformats.org/presentationml/2006/main}graphicFrame")
    assert len(frames) == 1


def test_series_child_order_is_schema_valid():
    _, root = _slide_and_chart()
    for ser in root.iter(f"{_C}ser"):
        tags = [etree.QName(ch).localname for ch in ser]
        assert tags.index("tx") < tags.index("spPr")
        assert tags.index("spPr") < tags.index("cat") < tags.index("val")


def test_series_colors_come_from_payload():
    _, root = _slide_and_chart()
    bar_chart = root.find(f".//{_C}barChart")
    fills = [el.get("val") for el in bar_chart.iter(f"{_A}srgbClr")]
    assert "94A3B8" in fills and "93C5FD" in fills
    line_chart = root.find(f".//{_C}lineChart")
    assert "EF4444" in [el.get("val") for el in line_chart.iter(f"{_A}srgbClr")]


def test_default_colors_when_not_specified():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    charts.add_combo_chart(slide, prs.part.package, ["a"],
                           [{"name": "b", "values": [1]}],
                           {"name": "l", "values": [2]},
                           x=0, y=0, cx=100, cy=100)
    part = next(p for p in prs.part.package.iter_parts()
                if "charts/chart" in str(p.partname))
    root = etree.fromstring(part.blob)
    fills = [el.get("val") for el in root.iter(f"{_A}srgbClr")]
    assert theme.CHART_BAR_COLORS[0] in fills
    assert theme.CHART_LINE_COLOR in fills


def test_gridlines_legend_and_axis_typography():
    _, root = _slide_and_chart()
    left_ax = root.findall(f".//{_C}valAx")[0]
    grid = left_ax.find(f"{_C}majorGridlines")
    assert grid is not None
    assert next(grid.iter(f"{_A}srgbClr")).get("val") == theme.CHART_GRID
    def_rpr = left_ax.find(f"{_C}txPr").find(f".//{_A}defRPr")
    assert def_rpr.get("sz") == "1000"
    assert def_rpr.find(f"{_A}latin").get("typeface") == theme.FONT
    legend = root.find(f".//{_C}legend")
    assert legend.find(f"{_C}legendPos").get("val") == "b"


def test_two_charts_get_distinct_partnames():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for _ in range(2):
        charts.add_combo_chart(slide, prs.part.package, ["a"],
                               [{"name": "b", "values": [1]}],
                               {"name": "l", "values": [2]},
                               x=0, y=0, cx=100, cy=100)
    names = sorted(str(p.partname) for p in prs.part.package.iter_parts()
                   if "charts/chart" in str(p.partname))
    assert names == ["/ppt/charts/chart1.xml", "/ppt/charts/chart2.xml"]
