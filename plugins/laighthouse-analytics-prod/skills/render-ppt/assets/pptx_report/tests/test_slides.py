from pptx import Presentation
from pptx.util import Emu

import slides as sl
import theme


def _prs():
    prs = Presentation()
    prs.slide_width = theme.SLIDE_W
    prs.slide_height = theme.SLIDE_H
    return prs


def _texts(slide):
    out = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            out.extend(p.text for p in shape.text_frame.paragraphs)
    return out


def test_cover_slide_title_and_period():
    prs = _prs()
    slide = sl.add_cover_slide(prs, "다형식품 MTD 보고서", "2026-05-01 ~ 2026-05-15")
    texts = _texts(slide)
    assert "다형식품 MTD 보고서" in texts
    assert any("2026-05-01" in t for t in texts)
    assert any("Engineered by Laighthouse AI" in t for t in texts)


def test_slide_background_is_page_tint():
    prs = _prs()
    slide = sl.add_cover_slide(prs, "t")
    assert str(slide.background.fill.fore_color.rgb) == theme.FILL_PAGE


def test_kpi_slide_cards_and_diff_colors():
    prs = _prs()
    slide = sl.add_kpi_slide(prs, "월 목표", [
        {"label": "월 매출 목표", "value": "₩ 250,000,000"},
        {"label": "ROAS", "value": "520%", "accent": "#3b82f6"},
        {"label": "증감", "value": "1", "diff": "▲ +2%", "diff_value": 2},
    ], page_no=2)
    texts = _texts(slide)
    assert "월 목표" in texts
    assert "₩ 250,000,000" in texts
    # find the accent-colored value run and the green diff run
    runs = [r for shape in slide.shapes if shape.has_text_frame
            for p in shape.text_frame.paragraphs for r in p.runs]
    accent_run = next(r for r in runs if r.text == "520%")
    assert str(accent_run.font.color.rgb) == "3B82F6"
    diff_run = next(r for r in runs if r.text == "▲ +2%")
    assert str(diff_run.font.color.rgb) == theme.GREEN


def test_table_slide_header_fill_and_no_default_style():
    prs = _prs()
    slide = sl.add_table_slide(prs, "캠페인별 성과", ["캠페인", "ROAS"],
                               [["브랜드검색", "727%"], ["합계", "573%"]], page_no=3)
    table = next(s for s in slide.shapes if s.has_table).table
    tbl_pr = table._tbl.tblPr
    assert tbl_pr.get("firstRow") == "0"
    assert tbl_pr.get("bandRow") == "0"
    header_cell = table.cell(0, 0)
    assert str(header_cell.fill.fore_color.rgb) == theme.FILL_HEADER
    # total row shaded + bold
    total_cell = table.cell(2, 0)
    assert str(total_cell.fill.fore_color.rgb) == theme.FILL_HEADER
    run = total_cell.text_frame.paragraphs[0].runs[0]
    assert run.font.bold is True
    assert run.font.name == theme.FONT


def test_table_slide_truncates_to_top_n_preserving_total():
    rows = [[f"항목{i}", str(i)] for i in range(30)] + [["합계", "999"]]
    prs = _prs()
    slide = sl.add_table_slide(prs, "표", ["이름", "값"], rows, max_rows=12)
    table = next(s for s in slide.shapes if s.has_table).table
    assert len(table.rows) == 1 + 12  # header + 12 body rows
    last_row_label = table.cell(12, 0).text_frame.text
    assert last_row_label == "합계"
    texts = _texts(slide)
    assert any("외 19행 생략" in t for t in texts)


def test_table_slide_rows_total_overrides_hidden_count():
    # source data was already truncated to 15 rows upstream; the caption
    # must report against the full dataset (rows_total), not the file
    rows = [[f"kw{i}", str(i)] for i in range(15)]
    prs = _prs()
    slide = sl.add_table_slide(prs, "키워드별 성과", ["키워드", "값"], rows,
                               max_rows=12, rows_total=1234)
    table = next(s for s in slide.shapes if s.has_table).table
    assert len(table.rows) == 1 + 12
    assert any("외 1222행 생략" in t for t in _texts(slide))


def test_table_slide_rows_total_with_preserved_total_row():
    rows = [[f"kw{i}", str(i)] for i in range(15)] + [["합계", "999"]]
    prs = _prs()
    slide = sl.add_table_slide(prs, "표", ["이름", "값"], rows,
                               max_rows=12, rows_total=1000)
    # 12 rows shown, one of them the 합계 row -> 11 body rows visible
    assert any("외 989행 생략" in t for t in _texts(slide))


def test_table_slide_no_truncation_note_when_small():
    prs = _prs()
    slide = sl.add_table_slide(prs, "표", ["a"], [["1"], ["2"]])
    assert not any("생략" in t for t in _texts(slide))


def test_table_slide_rich_and_bar_cells():
    prs = _prs()
    slide = sl.add_table_slide(prs, "표", ["지표", "달성률", "증감"], [
        ["매출",
         {"type": "bar", "pct": 53.0, "color": "3B82F6", "label": "53.0%"},
         {"text": "▲ +12.4%", "color": "#16a34a", "bold": True}],
    ])
    table = next(s for s in slide.shapes if s.has_table).table
    bar_run = table.cell(1, 1).text_frame.paragraphs[0].runs[0]
    assert bar_run.text == "53.0%"
    assert str(bar_run.font.color.rgb) == "3B82F6"
    rich_run = table.cell(1, 2).text_frame.paragraphs[0].runs[0]
    assert str(rich_run.font.color.rgb) == "16A34A"
    assert rich_run.font.bold is True


def test_table_columns_sized_by_content():
    rows = [["012_브랜드_공통_분유_견과류음료_견과류천국", "1,204,331", "12%"],
            ["짧은이름", "48,102", "7%"]]
    prs = _prs()
    slide = sl.add_table_slide(prs, "그룹별 성과", ["광고그룹", "노출", "CTR"], rows)
    table = next(s for s in slide.shapes if s.has_table).table
    widths = [table.columns[c].width for c in range(3)]
    assert widths[0] > widths[1] > widths[2]  # long-name column gets the most room
    assert abs(sum(widths) - int(theme.CONTENT_W)) <= 3  # still fills content width


def test_truncation_caption_sits_in_title_band():
    rows = [[f"kw{i}", str(i)] for i in range(30)]
    prs = _prs()
    slide = sl.add_table_slide(prs, "표", ["이름", "값"], rows, max_rows=12)
    note = next(s for s in slide.shapes if s.has_text_frame
                and "생략" in s.text_frame.text)
    assert note.top == theme.TITLE_TOP  # fixed band — cannot overlap grown rows
    assert "상위 12행 표시 · 외 18행 생략" in note.text_frame.text


def test_text_slide_multiline_in_card():
    prs = _prs()
    slide = sl.add_text_slide(prs, "Executive Summary", "첫 줄.\n둘째 줄.", page_no=4)
    texts = _texts(slide)
    assert "첫 줄." in texts
    assert "둘째 줄." in texts


def test_prettify_rounds_only_long_decimals():
    assert sl._prettify("55.88918788518661%") == "55.89%"
    assert sl._prettify("711,316.9333333333") == "711,316.93"
    assert sl._prettify("92335.39603960396%") == "92335.4%"
    assert sl._prettify("₩ 2,449.08 / 727% / 2026-05-15") == "₩ 2,449.08 / 727% / 2026-05-15"
    assert sl._prettify("합계") == "합계"


def test_table_cells_render_prettified_numbers():
    prs = _prs()
    slide = sl.add_table_slide(prs, "표", ["매체", "소진율"],
                               [["브랜드검색", "55.88918788518661%"]])
    table = next(s for s in slide.shapes if s.has_table).table
    assert table.cell(1, 1).text_frame.text == "55.89%"


def test_column_floor_prevents_vertical_headers():
    # many columns + one huge column must not squeeze "광고비" below its width
    rows = [["가나다라마바사아자차카타파하" * 2] + ["1"] * 7]
    headers = ["키워드", "노출", "클릭", "광고비", "CPC", "CPM", "구매건수", "ROAS"]
    prs = _prs()
    slide = sl.add_table_slide(prs, "표", headers, rows)
    table = next(s for s in slide.shapes if s.has_table).table
    widths = [table.columns[c].width for c in range(len(headers))]
    total = sum(widths)
    for header, width in zip(headers, widths):
        min_units = sl._display_width(header) + 2
        # each column's share must be at least its floor share of the total
        assert width / total >= min_units / (min_units * len(headers) + 40) * 0.5


def test_long_text_paginates_to_continuation_slide():
    long_line = "매출이 가장 많이 발생하는 카테고리를 분석 대상으로 선정하였습니다. " * 4
    body = "\n".join(["Overview"] + [long_line] * 10)
    prs = _prs()
    sl.add_text_slide(prs, "카테고리별 성과 분석", body)
    assert len(prs.slides) > 1
    titles = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                titles.append(shape.text_frame.text)
    assert any("카테고리별 성과 분석 (계속)" in t for t in titles)


def test_text_subheadings_render_bold():
    prs = _prs()
    slide = sl.add_text_slide(prs, "분석", "국내분유\n누적 매출 1위이며 성장했습니다.")
    runs = [r for shape in slide.shapes if shape.has_text_frame
            for p in shape.text_frame.paragraphs for r in p.runs]
    subhead = next(r for r in runs if r.text == "국내분유")
    assert subhead.font.bold is True
    body_run = next(r for r in runs if "누적 매출" in r.text)
    assert not body_run.font.bold


def test_divider_slide():
    prs = _prs()
    slide = sl.add_divider_slide(prs, "부록")
    assert "부록" in _texts(slide)
