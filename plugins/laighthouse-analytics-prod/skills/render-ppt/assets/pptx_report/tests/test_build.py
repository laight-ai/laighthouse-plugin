import json
import subprocess
import sys
import zipfile
from pathlib import Path

import pytest
from lxml import etree
from pptx import Presentation

import build

_DATA = {
    "title": "다형식품 MTD 보고서",
    "period": "2026-05-01 ~ 2026-05-15",
    "sections": [
        {"type": "kpi_cards", "cards": [{"label": "월 목표 매출", "value": "₩ 1,000"}]},
        {"type": "table", "heading": "목표 달성 현황", "headers": ["지표", "실적"],
         "rows": [["매출", "₩ 500"]]},
        {"type": "chart", "heading": "월별 광고 성과", "categories": ["1월", "2월"],
         "bar_series": [{"name": "광고비", "values": [1, 2]}],
         "line_series": {"name": "ROAS", "values": [3, 4]}},
        {"type": "text", "heading": "Executive Summary", "body": "요약."},
        {"type": "heading", "text": "부록"},
    ],
}


def _all_texts(prs):
    out = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                out.extend(p.text for p in shape.text_frame.paragraphs)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        out.append(cell.text_frame.text)
    return out


def test_build_presentation_one_slide_per_section_plus_cover():
    prs = build.build_presentation(_DATA)
    assert len(prs.slides) == 1 + len(_DATA["sections"])
    texts = _all_texts(prs)
    for expected in ("다형식품 MTD 보고서", "목표 달성 현황", "월별 광고 성과",
                     "요약.", "부록", "매출"):
        assert expected in texts


def test_standalone_heading_merges_into_headingless_follower():
    data = {"title": "t", "sections": [
        {"type": "heading", "text": "월 목표"},
        {"type": "kpi_cards", "cards": [{"label": "a", "value": "1"}]},
        {"type": "heading", "text": "부록"},  # trailing heading stays a divider
    ]}
    prs = build.build_presentation(data)
    assert len(prs.slides) == 3  # cover + merged kpi slide + divider
    kpi_texts = _all_texts(prs)
    assert "월 목표" in kpi_texts and "부록" in kpi_texts


def test_heading_before_headed_section_stays_divider():
    data = {"title": "t", "sections": [
        {"type": "heading", "text": "파트 1"},
        {"type": "text", "heading": "자체 제목", "body": "b"},
    ]}
    prs = build.build_presentation(data)
    assert len(prs.slides) == 3  # cover + divider + text slide


def test_build_presentation_rejects_unknown_type():
    with pytest.raises(ValueError, match="unknown section type"):
        build.build_presentation({"title": "t", "sections": [{"type": "nope"}]})


def test_cli_writes_pptx_with_wellformed_parts(tmp_path):
    data_path = tmp_path / "data.json"
    data_path.write_text(json.dumps(_DATA, ensure_ascii=False), encoding="utf-8")
    out_path = tmp_path / "out" / "report.pptx"
    subprocess.run(
        [sys.executable, str(Path(build.__file__)),
         "--data", str(data_path), "--out", str(out_path)],
        check=True,
    )
    assert out_path.exists()
    with zipfile.ZipFile(out_path) as z:
        assert any("charts/chart" in n for n in z.namelist())
        for name in z.namelist():
            if name.endswith(".xml") or name.endswith(".rels"):
                etree.fromstring(z.read(name))  # raises if malformed
    prs = Presentation(str(out_path))
    assert any("다형식품" in t for t in _all_texts(prs))
