"""python-pptx slide builders consumed by build.py.

Display-rounding note: MCP payloads carry raw float precision
(55.88918788518661% 등). _prettify() rounds any decimal run longer than two
places at the single point every rendered string passes through, so no
mapper needs its own rounding rule.

Each section-JSON object becomes one 16:9 slide. Unlike Word's flow layout,
slides give absolute positioning and real shapes, so the HTML report's card
UI (rounded corners, tinted page background, subtle borders) is reproduced
directly instead of approximated.
"""
import math
import re

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

import charts
import theme

_TOTAL_ROW_LABELS = {"합계", "총계"}
_NO_GRID_TABLE_STYLE = "{2D5ABB26-0587-4C30-8999-92F81FD0307C}"  # "No Style, No Grid"


# ── low-level helpers ──────────────────────────────────────────────────

def _rgb(hex_color):
    return RGBColor.from_string(hex_color.lstrip("#").upper())


_LONG_DECIMAL = re.compile(r"(\d+)\.(\d{3,})")


def _prettify(text):
    """Round runaway float precision for display: 55.88918788518661 → 55.89.
    Only decimal runs of 3+ places are touched, so pre-formatted values
    (2,449.08 / 727% / dates) pass through unchanged."""
    def _round(match):
        rounded = round(float(f"{match.group(1)}.{match.group(2)}"), 2)
        return f"{rounded:.2f}".rstrip("0").rstrip(".")
    return _LONG_DECIMAL.sub(_round, str(text))


def _style_run(run, size, color, bold=False):
    run.font.name = theme.FONT
    # Korean glyphs resolve through the eastAsia slot, not the latin one
    rpr = run._r.get_or_add_rPr()
    ea = rpr.find(qn("a:ea"))
    if ea is None:
        ea = rpr.makeelement(qn("a:ea"), {})
        rpr.append(ea)
    ea.set("typeface", theme.FONT)
    run.font.size = size
    run.font.color.rgb = _rgb(color)
    run.font.bold = bold
    return run


def _textbox(slide, x, y, w, h):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return box


def _add_line(tf, text, size, color, bold=False, first=False,
              align=PP_ALIGN.LEFT, space_before=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    if space_before is not None:
        p.space_before = space_before
    _style_run(p.add_run(), size, color, bold).text = _prettify(text)
    return p


def _card(slide, x, y, w, h, fill=theme.FILL_CARD, line=theme.BORDER):
    """Rounded-corner card — the .card CSS class made of an actual shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.adjustments[0] = 0.045
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill)
    shape.line.color.rgb = _rgb(line)
    shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    shape.text_frame.paragraphs[0].text = ""
    return shape


def _blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # fully blank layout
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(theme.FILL_PAGE)
    return slide


def _footer(slide, page_no=None):
    box = _textbox(slide, theme.MARGIN, theme.FOOTER_TOP,
                   theme.CONTENT_W, Emu(274320))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _style_run(p.add_run(), theme.SIZE_FOOTER, theme.TEXT_FAINT).text = \
        "Engineered by Laighthouse AI"
    if page_no is not None:
        num = _textbox(slide, theme.SLIDE_W - theme.MARGIN - Emu(457200),
                       theme.FOOTER_TOP, Emu(457200), Emu(274320))
        np = num.text_frame.paragraphs[0]
        np.alignment = PP_ALIGN.RIGHT
        _style_run(np.add_run(), theme.SIZE_FOOTER, theme.TEXT_FAINT).text = str(page_no)


def _slide_title(slide, text):
    """Accent bar + bold title — the .section-title analogue."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, theme.MARGIN, theme.TITLE_TOP,
        Emu(54864), theme.TITLE_H)  # 0.06 in wide
    bar.adjustments[0] = 0.5
    bar.fill.solid()
    bar.fill.fore_color.rgb = _rgb(theme.ACCENT)
    bar.line.fill.background()
    bar.shadow.inherit = False
    box = _textbox(slide, theme.MARGIN + Emu(137160), theme.TITLE_TOP,
                   theme.CONTENT_W - Emu(137160), theme.TITLE_H)
    box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    _add_line(box.text_frame, text, theme.SIZE_SLIDE_TITLE, theme.TEXT_STRONG,
              bold=True, first=True)


def _content_slide(prs, heading, page_no=None):
    slide = _blank_slide(prs)
    if heading:
        _slide_title(slide, heading)
    if page_no is None:
        page_no = len(prs.slides._sldIdLst)  # cover is 1, content follows
    _footer(slide, page_no)
    return slide


# ── slides ─────────────────────────────────────────────────────────────

def add_cover_slide(prs, title, period=None):
    slide = _blank_slide(prs)
    center_y = Emu(int(theme.SLIDE_H * 0.36))
    bar = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, theme.MARGIN, center_y - Emu(137160),
        Emu(731520), Emu(54864))  # short horizontal accent line
    bar.adjustments[0] = 0.5
    bar.fill.solid()
    bar.fill.fore_color.rgb = _rgb(theme.ACCENT)
    bar.line.fill.background()
    bar.shadow.inherit = False

    box = _textbox(slide, theme.MARGIN, center_y, theme.CONTENT_W, Emu(1371600))
    tf = box.text_frame
    _add_line(tf, title, theme.SIZE_COVER_TITLE, theme.TEXT_STRONG,
              bold=True, first=True)
    if period:
        p = _add_line(tf, f"📅 {period}", theme.SIZE_COVER_PERIOD, theme.TEXT_MUTED)
        p.space_before = Pt(10)
    _footer(slide)
    return slide


def add_divider_slide(prs, text, page_no=None):
    slide = _blank_slide(prs)
    box = _textbox(slide, theme.MARGIN, Emu(int(theme.SLIDE_H * 0.42)),
                   theme.CONTENT_W, Emu(914400))
    _add_line(box.text_frame, text, theme.SIZE_DIVIDER, theme.TEXT_STRONG,
              bold=True, first=True, align=PP_ALIGN.CENTER)
    _footer(slide, page_no if page_no is not None else len(prs.slides._sldIdLst))
    return slide


def add_kpi_slide(prs, heading, cards, page_no=None):
    slide = _content_slide(prs, heading, page_no)
    n = len(cards)
    gap = Emu(228600)  # 0.25 in
    card_w = Emu(int((theme.CONTENT_W - gap * (n - 1)) / n))
    card_h = Emu(1600200)  # 1.75 in
    y = theme.CONTENT_TOP + Emu(int((theme.CONTENT_H - card_h) * 0.35))
    for i, card in enumerate(cards):
        x = theme.MARGIN + Emu(int((card_w + gap) * i))
        shape = _card(slide, x, y, card_w, card_h)
        tf = shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_right = Emu(182880)
        _add_line(tf, card["label"], theme.SIZE_KPI_LABEL, theme.TEXT_MUTED,
                  first=True, align=PP_ALIGN.CENTER)
        accent = card.get("accent")
        value_p = _add_line(tf, card["value"], theme.SIZE_KPI_VALUE,
                            accent.lstrip("#") if accent else theme.TEXT_STRONG,
                            bold=True, align=PP_ALIGN.CENTER)
        value_p.space_before = Pt(6)
        if card.get("diff"):
            diff_p = _add_line(tf, card["diff"], theme.SIZE_KPI_DIFF,
                               theme.diff_color(card.get("diff_value")) or theme.GRAY,
                               bold=True, align=PP_ALIGN.CENTER)
            diff_p.space_before = Pt(4)
    return slide


def _is_bar_cell(value):
    return isinstance(value, dict) and value.get("type") == "bar"


def _is_rich_cell(value):
    return isinstance(value, dict) and "text" in value


def _truncate_rows(rows, limit):
    """상위 N개 규칙: keep the first rows, always preserving a trailing
    합계/총계 row, and report how many body rows were dropped."""
    if len(rows) <= limit:
        return rows, 0
    total_row = None
    body = rows
    first_cell = rows[-1][0] if rows[-1] else ""
    if _is_rich_cell(first_cell):
        first_cell = first_cell["text"]
    if str(first_cell) in _TOTAL_ROW_LABELS:
        total_row, body = rows[-1], rows[:-1]
    keep = limit - (1 if total_row is not None else 0)
    shown = body[:keep]
    hidden = len(body) - keep
    if total_row is not None:
        shown = shown + [total_row]
    return shown, hidden


def _display_width(value):
    """Approximate rendered width in half-width character units — CJK glyphs
    count double, so Korean campaign names get the width they need."""
    if _is_bar_cell(value):
        value = value.get("label", "")
    elif _is_rich_cell(value):
        value = value["text"]
    text = _prettify(value)  # widths must match what actually renders
    return sum(2 if ord(ch) > 0x2E80 else 1 for ch in text)


def _column_widths(headers, rows, total_width):
    """Distribute the table width by each column's longest content instead of
    equally — long name columns stop wrapping onto extra lines (which used to
    grow rows past the computed table height and overflow the slide)."""
    scores = []
    for c, header in enumerate(headers):
        header_width = _display_width(header)
        longest = header_width
        for row in rows:
            if c < len(row):
                longest = max(longest, _display_width(row[c]))
        # floor at the header's own width so short columns never wrap their
        # header vertically ("광고비" one glyph per line), cap runaway cells
        scores.append(min(max(longest, header_width + 2, 8), 40))
    total_score = sum(scores)
    widths = [Emu(int(total_width * s / total_score)) for s in scores]
    widths[-1] = Emu(int(total_width) - sum(int(w) for w in widths[:-1]))
    return widths


def _fill_cell(cell, text, size, color, bold=False, fill=None):
    cell.margin_left = cell.margin_right = Emu(109728)
    cell.margin_top = cell.margin_bottom = Emu(45720)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if fill:
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb(fill)
    else:
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb(theme.FILL_CARD)
    _add_line(cell.text_frame, text, size, color, bold, first=True)
    # bottom hairline — the HTML row's border-bottom
    tc_pr = cell._tc.get_or_add_tcPr()
    ln_b = tc_pr.makeelement(qn("a:lnB"), {"w": "6350", "cap": "flat"})
    fill_el = ln_b.makeelement(qn("a:solidFill"), {})
    clr = fill_el.makeelement(qn("a:srgbClr"), {"val": theme.BORDER_SOFT})
    fill_el.append(clr)
    ln_b.append(fill_el)
    tc_pr.append(ln_b)


def add_table_slide(prs, heading, headers, rows, page_no=None,
                    max_rows=theme.MAX_TABLE_ROWS, rows_total=None):
    """rows_total: original body-row count when the data was already
    truncated upstream (SKILL.md's top-15 저장 규칙) — keeps the
    "외 n행 생략" caption honest about the full dataset."""
    slide = _content_slide(prs, heading, page_no)
    shown, hidden = _truncate_rows(rows, max_rows)

    def _is_total(row):
        first = row[0] if row else ""
        if _is_rich_cell(first):
            first = first["text"]
        return str(first) in _TOTAL_ROW_LABELS

    body_shown = sum(1 for row in shown if not _is_total(row))
    if rows_total is not None:
        hidden = max(rows_total - body_shown, 0)

    header_h = Emu(320040)
    row_h = Emu(292608)
    table_h = Emu(int(header_h + row_h * len(shown)))
    gf = slide.shapes.add_table(1 + len(shown), len(headers),
                                theme.MARGIN, theme.CONTENT_TOP,
                                theme.CONTENT_W, table_h)
    table = gf.table
    # drop PowerPoint's default banded blue theme entirely
    tbl_pr = table._tbl.tblPr
    tbl_pr.set("firstRow", "0")
    tbl_pr.set("bandRow", "0")
    style_el = tbl_pr.find(qn("a:tableStyleId"))
    if style_el is None:
        style_el = tbl_pr.makeelement(qn("a:tableStyleId"), {})
        tbl_pr.append(style_el)
    style_el.text = _NO_GRID_TABLE_STYLE

    table.rows[0].height = header_h
    for r in range(1, 1 + len(shown)):
        table.rows[r].height = row_h

    for c, width in enumerate(_column_widths(headers, shown, theme.CONTENT_W)):
        table.columns[c].width = width

    for c, header in enumerate(headers):
        _fill_cell(table.cell(0, c), header, theme.SIZE_TH, theme.TEXT_TH,
                   bold=True, fill=theme.FILL_HEADER)

    for r, row in enumerate(shown, start=1):
        first_cell = row[0] if row else ""
        if _is_rich_cell(first_cell):
            first_cell = first_cell["text"]
        is_total = str(first_cell) in _TOTAL_ROW_LABELS
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            fill = theme.FILL_HEADER if is_total else None
            if _is_bar_cell(value):
                _fill_cell(cell, value.get("label", f"{value['pct']}%"),
                           theme.SIZE_TD, value.get("color", theme.ACCENT),
                           bold=True, fill=fill)
            elif _is_rich_cell(value):
                _fill_cell(cell, value["text"], theme.SIZE_TD,
                           value.get("color", theme.TEXT_TABLE),
                           bold=value.get("bold", is_total), fill=fill)
            else:
                _fill_cell(cell, value, theme.SIZE_TD, theme.TEXT_TABLE,
                           bold=is_total, fill=fill)

    if hidden:
        # anchored to the fixed title band, not below the table — wrapped
        # rows can grow the table past its computed height, and a caption
        # positioned by table_h then overlaps the rows
        note = _textbox(slide, theme.MARGIN, theme.TITLE_TOP,
                        theme.CONTENT_W, theme.TITLE_H)
        note.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        _add_line(note.text_frame, f"상위 {body_shown}행 표시 · 외 {hidden}행 생략",
                  theme.SIZE_TABLE_NOTE, theme.TEXT_FAINT, first=True,
                  align=PP_ALIGN.RIGHT)
    return slide


# pagination model for text slides: estimated wrapped lines per card, from
# the card's inner width at SIZE_BODY (display units: CJK counts double)
_BODY_UNITS_PER_LINE = 130
_BODY_LINES_PER_CARD = 14


def _is_subheading(line):
    """Short label lines between paragraphs (Overview / 국내분유 / 커피)
    render as bold subheadings — sentence lines end with a stop."""
    text = line.strip()
    return bool(text) and _display_width(text) <= 30 and not text.endswith((".", "!", "?", "%"))


def _paginate_body(body):
    chunks, current, used = [], [], 0
    for line in str(body).split("\n"):
        est = max(1, math.ceil(_display_width(line) / _BODY_UNITS_PER_LINE))
        if current and used + est > _BODY_LINES_PER_CARD:
            chunks.append(current)
            current, used = [], 0
        current.append(line)
        used += est
    chunks.append(current)
    return chunks


def add_text_slide(prs, heading, body, page_no=None):
    """Long analysis text flows onto continuation slides ("(계속)") instead
    of overflowing the card past the slide bottom."""
    first_slide = None
    for idx, chunk in enumerate(_paginate_body(body)):
        title = heading if idx == 0 else (f"{heading} (계속)" if heading else None)
        slide = _content_slide(prs, title, page_no if idx == 0 else None)
        first_slide = first_slide or slide
        card = _card(slide, theme.MARGIN, theme.CONTENT_TOP,
                     theme.CONTENT_W, theme.CONTENT_H)
        tf = card.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_left = tf.margin_right = Emu(274320)
        tf.margin_top = tf.margin_bottom = Emu(228600)
        for i, line in enumerate(chunk):
            if _is_subheading(line):
                p = _add_line(tf, line, theme.SIZE_BODY_SUBHEAD, theme.TEXT_STRONG,
                              bold=True, first=(i == 0))
                p.space_before = Pt(12)
                p.space_after = Pt(2)
            else:
                p = _add_line(tf, line, theme.SIZE_BODY, theme.TEXT_TABLE,
                              first=(i == 0))
                p.space_after = Pt(6)
            p.line_spacing = 1.35
    return first_slide


def add_chart_slide(prs, heading, categories, bar_series, line_series,
                    page_no=None):
    slide = _content_slide(prs, heading, page_no)
    _card(slide, theme.MARGIN, theme.CONTENT_TOP, theme.CONTENT_W, theme.CONTENT_H)
    pad = Emu(182880)
    charts.add_combo_chart(
        slide, prs.part.package, categories, bar_series, line_series,
        x=theme.MARGIN + pad, y=theme.CONTENT_TOP + pad,
        cx=theme.CONTENT_W - 2 * pad, cy=theme.CONTENT_H - 2 * pad,
    )
    return slide
