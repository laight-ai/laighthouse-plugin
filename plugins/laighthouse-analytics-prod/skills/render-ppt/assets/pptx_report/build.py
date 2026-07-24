"""CLI entrypoint: JSON section list -> .pptx deck.

Invoked directly as a script (not via -m), so it adds its own directory to
sys.path to resolve the sibling `slides`/`charts`/`theme` modules regardless
of the caller's working directory.

Section-JSON contract is shared with the mapping layer (section_mapping.py):
each object in `sections` has a `type` of kpi_cards / table / chart / text /
heading and becomes exactly one slide, in order, after the cover slide.
"""
import argparse
import json
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))

from pptx import Presentation  # noqa: E402

import slides as sl  # noqa: E402
import theme  # noqa: E402


def _merge_headings(sections):
    """A standalone heading followed by a heading-less section becomes that
    section's slide title instead of wasting a divider slide."""
    merged = []
    pending = None
    for section in sections:
        if section["type"] == "heading":
            if pending is not None:
                merged.append(pending)
            pending = section
            continue
        if pending is not None:
            if not section.get("heading"):
                section = {**section, "heading": pending["text"]}
            else:
                merged.append(pending)
            pending = None
        merged.append(section)
    if pending is not None:
        merged.append(pending)
    return merged


def _fill_analysis_slots(sections, analysis):
    """Replace map_report.py's analysis_slot placeholders with the LLM-written
    text sections; a slot with no matching analysis entry degrades to the
    데이터 준비 중 rule instead of failing the whole deck."""
    filled = []
    for section in sections:
        if section.get("type") != "analysis_slot":
            filled.append(section)
            continue
        entry = (analysis or {}).get(section["key"])
        if entry and entry.get("body"):
            filled.append({"type": "text",
                           "heading": entry.get("heading", section.get("heading")),
                           "body": entry["body"]})
        else:
            filled.append({"type": "text", "heading": section.get("heading"),
                           "body": "데이터 준비 중"})
    return filled


def build_presentation(data, analysis=None):
    prs = Presentation()
    prs.slide_width = theme.SLIDE_W
    prs.slide_height = theme.SLIDE_H
    sl.add_cover_slide(prs, data["title"], data.get("period"))
    sections = _fill_analysis_slots(data["sections"], analysis)
    for i, section in enumerate(_merge_headings(sections)):
        stype = section["type"]
        page_no = i + 2  # cover slide is page 1
        if stype == "heading":
            sl.add_divider_slide(prs, section["text"], page_no)
        elif stype == "kpi_cards":
            sl.add_kpi_slide(prs, section.get("heading"), section["cards"], page_no)
        elif stype == "table":
            sl.add_table_slide(prs, section.get("heading"), section["headers"],
                               section["rows"], page_no,
                               rows_total=section.get("rows_total"))
        elif stype == "chart":
            sl.add_chart_slide(prs, section.get("heading"), section["categories"],
                               section["bar_series"], section["line_series"], page_no)
        elif stype == "text":
            sl.add_text_slide(prs, section.get("heading"), section["body"], page_no)
        else:
            raise ValueError(f"unknown section type: {stype}")
    return prs


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--data", required=True, help="path to the section-list JSON file")
    parser.add_argument("--out", required=True, help="path to write the .pptx file to")
    parser.add_argument("--analysis", default=None,
                        help="optional JSON filling map_report.py analysis_slot "
                             "placeholders: {\"section3\": {\"heading\": ..., \"body\": ...}}")
    args = parser.parse_args()

    data = json.loads(Path(args.data).read_text(encoding="utf-8"))
    analysis = None
    if args.analysis:
        analysis = json.loads(Path(args.analysis).read_text(encoding="utf-8"))
    prs = build_presentation(data, analysis)

    out_path = Path(args.out).expanduser()
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))


if __name__ == "__main__":
    main()
